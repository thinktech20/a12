"""Extract all promo input docs + templates to plain text/markdown for review."""

from pathlib import Path
import openpyxl
from pptx import Presentation
import fitz  # PyMuPDF

ROOT = Path("/home/u560060992/dbx/ai-arch/promo")
IN_DIR = ROOT / "input-docs"
OUT_DIR = ROOT / "extracted"
OUT_DIR.mkdir(exist_ok=True)


def extract_xlsx(path: Path) -> str:
    """Dump every sheet as Markdown tables."""
    wb = openpyxl.load_workbook(path, data_only=True)
    out = [f"# {path.name}\n"]
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        out.append(f"\n## Sheet: {sheet_name}\n")
        rows = list(ws.iter_rows(values_only=True))
        # trim trailing all-None rows
        while rows and all(c is None for c in rows[-1]):
            rows.pop()
        if not rows:
            out.append("_(empty sheet)_\n")
            continue
        # determine max col with any value
        max_col = max((i for r in rows for i, c in enumerate(r) if c is not None), default=0) + 1
        for r_idx, row in enumerate(rows):
            cells = ["" if c is None else str(c).replace("\n", " / ").replace("|", "\\|") for c in row[:max_col]]
            out.append(f"| {' | '.join(cells)} |")
            if r_idx == 0:
                out.append("| " + " | ".join(["---"] * max_col) + " |")
        out.append("")
    return "\n".join(out)


def extract_pptx(path: Path) -> str:
    """Dump each slide as a section with title + bullets/text."""
    prs = Presentation(path)
    out = [f"# {path.name}\n"]
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"\n## Slide {i}\n")
        # try to find title
        title = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text_frame.text.strip()
        if title:
            out.append(f"**Title:** {title}\n")
        # gather all text frames (including tables)
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    out.append(txt)
                    out.append("")
            if shape.has_table:
                tbl = shape.table
                out.append("\n_Table:_")
                for row in tbl.rows:
                    cells = [c.text.strip().replace("\n", " / ").replace("|", "\\|") for c in row.cells]
                    out.append(f"| {' | '.join(cells)} |")
                out.append("")
        # speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                out.append(f"\n_Speaker notes:_ {notes}")
    return "\n".join(out)


def extract_pdf(path: Path) -> str:
    """Dump each page as plain text."""
    doc = fitz.open(path)
    out = [f"# {path.name}\n"]
    for i in range(doc.page_count):
        out.append(f"\n## Page {i+1}\n")
        out.append(doc.load_page(i).get_text().strip())
    doc.close()
    return "\n".join(out)


HANDLERS = {".xlsx": extract_xlsx, ".pptx": extract_pptx, ".pdf": extract_pdf}


def main():
    inputs = []
    for p in IN_DIR.rglob("*"):
        if p.is_file() and p.suffix.lower() in HANDLERS and ":Zone.Identifier" not in p.name:
            inputs.append(p)

    for src in sorted(inputs):
        rel = src.relative_to(IN_DIR)
        out_name = rel.with_suffix(".md").as_posix().replace("/", "__")
        dest = OUT_DIR / out_name
        try:
            text = HANDLERS[src.suffix.lower()](src)
            dest.write_text(text, encoding="utf-8")
            print(f"OK   {src.name}  ->  extracted/{out_name}  ({len(text):,} chars)")
        except Exception as e:
            print(f"FAIL {src.name}: {type(e).__name__}: {e}")


if __name__ == "__main__":
    main()
